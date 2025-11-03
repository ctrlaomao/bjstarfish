#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
import sys

def analyze_template(template_path):
    """分析PPT模板的样式特征"""
    prs = Presentation(template_path)
    
    print("=" * 60)
    print("模板分析报告：{}".format(template_path))
    print("=" * 60)
    
    # 1. 基本信息
    print("\n【基本信息】")
    print("幻灯片尺寸：宽 {} x 高 {}".format(prs.slide_width, prs.slide_height))
    print("总幻灯片数：{}".format(len(prs.slides)))
    print("可用布局数：{}".format(len(prs.slide_layouts)))
    
    # 2. 母版和布局
    print("\n【母版和布局】")
    for idx, layout in enumerate(prs.slide_layouts):
        print("布局 {}: {} (包含 {} 个占位符)".format(
            idx, layout.name, len(layout.placeholders)))
    
    # 3. 分析前5张幻灯片的样式
    print("\n【幻灯片内容分析（前5张）】")
    slides_list = list(prs.slides)
    for idx, slide in enumerate(slides_list[:5]):
        print("\n--- 幻灯片 {} ---".format(idx + 1))
        print("使用布局：{}".format(slide.slide_layout.name))
        
        # 分析文本框和占位符
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_preview = shape.text[:50].replace('\n', ' ')
                print("  文本: \"{}\"".format(text_preview))
                
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        if len(para.runs) > 0:
                            run = para.runs[0]
                            print("    字体: {}, 大小: {}, 加粗: {}".format(
                                run.font.name,
                                run.font.size,
                                run.font.bold
                            ))
                            if run.font.color and hasattr(run.font.color, 'rgb'):
                                print("    颜色: RGB{}".format(run.font.color.rgb))
                            break
    
    # 4. 检查背景
    print("\n【背景样式】")
    if len(prs.slides) > 0:
        first_slide = prs.slides[0]
        if hasattr(first_slide.background, 'fill'):
            print("第一张幻灯片背景类型：{}".format(first_slide.background.fill.type))
    
    print("\n" + "=" * 60)
    print("分析完成！")
    print("=" * 60)

if __name__ == "__main__":
    template_path = "/workspace/海星育数字化系统-0627.pptx"
    try:
        analyze_template(template_path)
    except Exception as e:
        print("分析出错：{}".format(str(e)))
        import traceback
        traceback.print_exc()
