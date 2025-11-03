#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
import copy

# 中文字体
CHINESE_FONT = "WenQuanYi Zen Hei"
# 模板主色
MAIN_COLOR = RGBColor(0, 70, 160)

def set_font(paragraph, font_name):
    """设置段落字体，包括东亚字体"""
    try:
        for run in paragraph.runs:
            run.font.name = font_name
            if hasattr(run.font, '_element'):
                rPr = run.font._element.get_or_add_rPr()
                ea = OxmlElement('a:ea')
                ea.set('typeface', font_name)
                rPr.append(ea)
                cs = OxmlElement('a:cs')
                cs.set('typeface', font_name)
                rPr.append(cs)
    except:
        pass

def extract_content(content_ppt_path):
    """提取内容PPT的所有数据"""
    print("正在读取内容文件：{}".format(content_ppt_path))
    prs = Presentation(content_ppt_path)
    
    slides_data = []
    slides_list = list(prs.slides)
    
    for idx, slide in enumerate(slides_list):
        slide_info = {
            'index': idx + 1,
            'texts': []
        }
        
        # 提取所有文本块
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_info['texts'].append(shape.text.strip())
        
        slides_data.append(slide_info)
        print("  提取幻灯片 {}: {} 个文本块".format(idx + 1, len(slide_info['texts'])))
    
    print("✓ 共提取 {} 张幻灯片的内容\n".format(len(slides_data)))
    return slides_data

def create_merged_ppt(template_path, slides_data, output_path):
    """创建合并后的PPT"""
    print("正在加载模板：{}".format(template_path))
    
    # 加载模板
    prs = Presentation(template_path)
    
    # 获取布局
    title_layout = prs.slide_layouts[12]  # 1_标题幻灯片
    section_layout = prs.slide_layouts[13]  # 1_节标题
    blank_layout = prs.slide_layouts[6]  # 空白
    
    print("正在清空模板幻灯片...")
    # 清空所有现有幻灯片
    slide_ids = [slide for slide in prs.slides._sldIdLst]
    for slide_id in slide_ids:
        prs.slides._sldIdLst.remove(slide_id)
    
    print("开始融合内容...\n")
    
    # 遍历每张内容幻灯片
    for slide_data in slides_data:
        idx = slide_data['index']
        texts = slide_data['texts']
        
        if not texts:
            continue
        
        print("生成幻灯片 {}/{}".format(idx, len(slides_data)))
        
        # 判断幻灯片类型
        first_text = texts[0] if texts else ""
        
        # 标题页（第一张）
        if idx == 1:
            add_title_slide(prs, title_layout, texts)
            print("  类型：标题页")
        
        # 节标题（一、二、三等开头）
        elif first_text.startswith(("一、", "二、", "三、", "四、", "五、", "六、")):
            section_num = {"一、": "01", "二、": "02", "三、": "03", 
                          "四、": "04", "五、": "05", "六、": "06"}.get(first_text[:2], "00")
            section_title = first_text[2:]
            
            # 如果只有标题没有内容，用节标题布局
            if len(texts) == 1:
                add_section_slide(prs, section_layout, section_num, section_title)
                print("  类型：节标题 - {}".format(section_title))
            else:
                # 有内容，用内容页
                add_content_slide(prs, blank_layout, first_text, texts[1:])
                print("  类型：内容页 - {}".format(section_title))
        
        # 结束页
        elif "谢谢" in first_text or "Q & A" in first_text:
            add_end_slide(prs, blank_layout, texts)
            print("  类型：结束页")
        
        # 普通内容页
        else:
            add_content_slide(prs, blank_layout, first_text, texts[1:] if len(texts) > 1 else [])
            print("  类型：内容页 - {}".format(first_text[:20]))
    
    # 保存
    print("\n正在保存：{}".format(output_path))
    prs.save(output_path)
    
    print("\n" + "=" * 70)
    print("✅ PPT合并完成！")
    print("📄 文件：{}".format(output_path))
    print("📊 共生成：{} 张幻灯片".format(len(slides_data)))
    print("=" * 70)

def add_title_slide(prs, layout, texts):
    """添加标题页"""
    slide = prs.slides.add_slide(layout)
    
    title = texts[0] if len(texts) > 0 else "项目汇报"
    subtitle = texts[1] if len(texts) > 1 else ""
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.name = CHINESE_FONT
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = MAIN_COLOR
    title_para.alignment = PP_ALIGN.CENTER
    set_font(title_para, CHINESE_FONT)
    
    # 副标题
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11.33), Inches(0.8))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.name = CHINESE_FONT
        subtitle_para.font.size = Pt(28)
        subtitle_para.font.color.rgb = RGBColor(89, 89, 89)
        subtitle_para.alignment = PP_ALIGN.CENTER
        set_font(subtitle_para, CHINESE_FONT)

def add_section_slide(prs, layout, number, title):
    """添加节标题页"""
    slide = prs.slides.add_slide(layout)
    
    # 清空占位符
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            for para in shape.text_frame.paragraphs:
                para.clear()
    
    # 数字
    num_box = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(9.33), Inches(1.5))
    num_frame = num_box.text_frame
    num_frame.text = number
    num_para = num_frame.paragraphs[0]
    num_para.font.name = CHINESE_FONT
    num_para.font.size = Pt(100)
    num_para.font.bold = True
    num_para.font.color.rgb = MAIN_COLOR
    num_para.alignment = PP_ALIGN.CENTER
    set_font(num_para, CHINESE_FONT)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(2), Inches(3.8), Inches(9.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.name = CHINESE_FONT
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = MAIN_COLOR
    title_para.alignment = PP_ALIGN.CENTER
    set_font(title_para, CHINESE_FONT)

def add_content_slide(prs, layout, title, content_texts):
    """添加内容页 - 智能调整字体大小"""
    slide = prs.slides.add_slide(layout)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.name = CHINESE_FONT
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = MAIN_COLOR
    set_font(title_para, CHINESE_FONT)
    
    # 内容区域
    if content_texts:
        # 合并所有内容文本
        all_content = "\n".join(content_texts)
        lines = all_content.split('\n')
        
        # 根据内容量动态调整字体大小
        total_lines = len(lines)
        if total_lines > 50:
            base_font_size = 11
            title_font_size = 13
        elif total_lines > 35:
            base_font_size = 12
            title_font_size = 14
        elif total_lines > 25:
            base_font_size = 13
            title_font_size = 15
        else:
            base_font_size = 14
            title_font_size = 16
        
        # 内容文本框
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.33), Inches(5.8))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        for i, line in enumerate(lines):
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            
            p.text = line
            p.font.name = CHINESE_FONT
            p.font.size = Pt(base_font_size)
            p.space_after = Pt(3)
            
            # 设置缩进和层级
            if line.startswith('•'):
                p.level = 0
            elif line.startswith('   •'):
                p.level = 1
                p.font.size = Pt(base_font_size - 1)
            elif line.startswith('   '):
                p.level = 1
                p.font.size = Pt(base_font_size - 1)
            
            # 加粗小标题
            if line and not line.startswith(' ') and not line.startswith('•') and \
               len(line) < 30 and ':' not in line and '：' not in line and \
               not line.startswith(('第', '已', '当', '下', '验收', '潜在')):
                p.font.bold = True
                p.font.size = Pt(title_font_size)
                p.font.color.rgb = MAIN_COLOR
            
            set_font(p, CHINESE_FONT)

def add_end_slide(prs, layout, texts):
    """添加结束页"""
    slide = prs.slides.add_slide(layout)
    
    # 主文本
    text = "\n\n".join(texts)
    
    end_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9.33), Inches(2))
    end_frame = end_box.text_frame
    end_frame.text = text
    
    for para in end_frame.paragraphs:
        para.font.name = CHINESE_FONT
        para.font.size = Pt(52)
        para.font.bold = True
        para.font.color.rgb = MAIN_COLOR
        para.alignment = PP_ALIGN.CENTER
        set_font(para, CHINESE_FONT)

if __name__ == "__main__":
    # 1. 提取内容
    content_path = "/workspace/report_ppt/内容.pptx"
    slides_data = extract_content(content_path)
    
    # 2. 创建合并PPT
    template_path = "/workspace/report_ppt/模板.pptx"
    output_path = "/workspace/report_ppt/合并后的汇报PPT.pptx"
    
    create_merged_ppt(template_path, slides_data, output_path)
