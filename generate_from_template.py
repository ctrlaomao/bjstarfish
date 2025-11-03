#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
import copy

# ä¸­æ–‡å­—ä½“
CHINESE_FONT = "WenQuanYi Zen Hei"
# æ¨¡æ¿ä¸»è‰²
MAIN_COLOR = RGBColor(0, 70, 160)

def set_font(paragraph, font_name):
    """è®¾ç½®æ®µè½å­—ä½“ï¼ŒåŒ…æ‹¬ä¸œäºšå­—ä½“"""
    try:
        for run in paragraph.runs:
            run.font.name = font_name
            if hasattr(run.font, '_element'):
                rPr = run.font._element.get_or_add_rPr()
                ea = OxmlElement('a:ea')
                ea.set('typeface', font_name)
                rPr.append(ea)
                cs = OxmlElement('a:cs')
                cs.set('typeface', font_name)
                rPr.append(cs)
    except:
        pass

def create_presentation_from_template():
    """åŸºäºæ¨¡æ¿åˆ›å»ºæ–°çš„æ±‡æŠ¥PPT"""
    
    # åŠ è½½æ¨¡æ¿
    template_path = "/workspace/æµ·æ˜Ÿè‚²æ•°å­—åŒ–ç³»ç»Ÿ-0627.pptx"
    prs = Presentation(template_path)
    
    # è·å–å¸ƒå±€
    title_layout = prs.slide_layouts[12]  # 1_æ ‡é¢˜å¹»ç¯ç‰‡
    section_layout = prs.slide_layouts[13]  # 1_èŠ‚æ ‡é¢˜
    content_layout = prs.slide_layouts[14]  # 1_ä¸¤æ å†…å®¹
    blank_layout = prs.slide_layouts[6]  # ç©ºç™½
    
    # æ¸…ç©ºæ‰€æœ‰å¹»ç¯ç‰‡
    slide_ids = [slide for slide in prs.slides._sldIdLst]
    for slide_id in slide_ids:
        prs.slides._sldIdLst.remove(slide_id)
    
    print("å¼€å§‹ç”Ÿæˆæ–°çš„æ±‡æŠ¥PPT...")
    
    # æ ‡é¢˜é¡µ
    slide = prs.slides.add_slide(title_layout)
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            for para in shape.text_frame.paragraphs:
                if len(para.runs) > 0:
                    # æ¸…ç©ºæ‰€æœ‰å ä½ç¬¦æ–‡æœ¬
                    para.clear()
    
    # æ‰‹åŠ¨æ·»åŠ æ ‡é¢˜å†…å®¹
    add_custom_title_slide(prs, slide, "å‘¨ä¸‰é¡¹ç›®æ±‡æŠ¥", "èƒœè¾¾è®¯é¡¹ç›®å›¢é˜Ÿ", "2023å¹´11æœˆ")
    
    # ç¬¬ä¸€éƒ¨åˆ†ï¼šé¡¹ç›®è§„åˆ’
    add_section_title(prs, section_layout, "01", "é¡¹ç›®è§„åˆ’")
    add_content_slide(prs, blank_layout, "ä¸€ã€é¡¹ç›®è§„åˆ’", [
        "é¡¹ç›®ç›®æ ‡",
        "â€¢ æ„å»ºé«˜æ•ˆã€å¯é çš„ä¸šåŠ¡ç³»ç»Ÿ",
        "â€¢ æå‡ä¸šåŠ¡å¤„ç†æ•ˆç‡å’Œç”¨æˆ·ä½“éªŒ",
        "â€¢ å®ç°æ•°æ®é©±åŠ¨çš„å†³ç­–æ”¯æŒ",
        "â€¢ ä¿éšœç³»ç»Ÿå®‰å…¨æ€§å’Œç¨³å®šæ€§",
        "",
        "é¡¹ç›®èŒƒå›´",
        "â€¢ ç³»ç»ŸåŠŸèƒ½æ¨¡å—å¼€å‘",
        "â€¢ æ•°æ®åº“è®¾è®¡ä¸å®ç°",
        "â€¢ æ¥å£å¼€å‘ä¸é›†æˆ",
        "â€¢ ç³»ç»Ÿæµ‹è¯•ä¸éƒ¨ç½²",
        "â€¢ ç”¨æˆ·åŸ¹è®­ä¸æ–‡æ¡£ç¼–åˆ¶",
        "",
        "é¡¹ç›®å‘¨æœŸ",
        "â€¢ å¯åŠ¨æ—¶é—´ï¼š2023å¹´4æœˆ",
        "â€¢ å¼€å‘å‘¨æœŸï¼š2023å¹´4æœˆ - 10æœˆ",
        "â€¢ å½“å‰çŠ¶æ€ï¼šå‡†å¤‡éªŒæ”¶é˜¶æ®µ"
    ])
    
    # ç¬¬äºŒéƒ¨åˆ†ï¼šä¸šåŠ¡æµç¨‹
    add_section_title(prs, section_layout, "02", "ä¸šåŠ¡æµç¨‹")
    add_content_slide(prs, blank_layout, "äºŒã€ä¸šåŠ¡æµç¨‹", [
        "æ ¸å¿ƒä¸šåŠ¡æµç¨‹",
        "",
        "1. éœ€æ±‚æ”¶é›†ä¸åˆ†æ",
        "   â€¢ ç”¨æˆ·éœ€æ±‚è°ƒç ”",
        "   â€¢ ä¸šåŠ¡æµç¨‹æ¢³ç†",
        "   â€¢ éœ€æ±‚æ–‡æ¡£ç¼–å†™",
        "   â€¢ éœ€æ±‚è¯„å®¡ä¸ç¡®è®¤",
        "",
        "2. ç³»ç»Ÿè®¾è®¡",
        "   â€¢ æ€»ä½“æ¶æ„è®¾è®¡",
        "   â€¢ æ•°æ®åº“è¯¦ç»†è®¾è®¡",
        "   â€¢ æ¥å£è®¾è®¡ä¸è§„èŒƒ",
        "   â€¢ å®‰å…¨æœºåˆ¶è®¾è®¡",
        "",
        "3. å¼€å‘å®ç°",
        "   â€¢ æ¨¡å—åŒ–å¼€å‘",
        "   â€¢ ä»£ç å®¡æŸ¥",
        "   â€¢ å•å…ƒæµ‹è¯•",
        "   â€¢ é›†æˆæµ‹è¯•"
    ])
    
    add_content_slide(prs, blank_layout, "äºŒã€ä¸šåŠ¡æµç¨‹ï¼ˆç»­ï¼‰", [
        "4. ç³»ç»Ÿæµ‹è¯•",
        "   â€¢ åŠŸèƒ½æµ‹è¯•",
        "   â€¢ æ€§èƒ½æµ‹è¯•",
        "   â€¢ å®‰å…¨æµ‹è¯•",
        "   â€¢ ç”¨æˆ·éªŒæ”¶æµ‹è¯•ï¼ˆUATï¼‰",
        "",
        "5. éƒ¨ç½²ä¸Šçº¿",
        "   â€¢ ç”Ÿäº§ç¯å¢ƒå‡†å¤‡",
        "   â€¢ ç³»ç»Ÿéƒ¨ç½²å®æ–½",
        "   â€¢ æ•°æ®è¿ç§»ä¸éªŒè¯",
        "   â€¢ åˆ‡æ¢æ–¹æ¡ˆæ‰§è¡Œ",
        "",
        "6. è¿ç»´æ”¯æŒ",
        "   â€¢ 7Ã—24å°æ—¶ç³»ç»Ÿç›‘æ§",
        "   â€¢ å¿«é€Ÿé—®é¢˜å“åº”æœºåˆ¶",
        "   â€¢ æŒç»­ä¼˜åŒ–ä¸æ”¹è¿›",
        "   â€¢ å®šæœŸå·¡æ£€ä¸ç»´æŠ¤"
    ])
    
    # ç¬¬ä¸‰éƒ¨åˆ†ï¼šç³»ç»Ÿæ¼”ç¤º
    add_section_title(prs, section_layout, "03", "ç³»ç»Ÿæ¼”ç¤º")
    add_content_slide(prs, blank_layout, "ä¸‰ã€ç³»ç»Ÿæ¼”ç¤º", [
        "ç³»ç»Ÿæ¶æ„",
        "â€¢ å‰ç«¯å±•ç¤ºå±‚ï¼šå“åº”å¼Webç•Œé¢",
        "â€¢ ä¸šåŠ¡é€»è¾‘å±‚ï¼šRESTful APIæœåŠ¡",
        "â€¢ æ•°æ®è®¿é—®å±‚ï¼šORMæ¡†æ¶å°è£…",
        "â€¢ æ•°æ®å­˜å‚¨å±‚ï¼šå…³ç³»å‹æ•°æ®åº“ + ç¼“å­˜",
        "â€¢ åŸºç¡€è®¾æ–½å±‚ï¼šäº‘æœåŠ¡æ”¯æŒ",
        "",
        "æ ¸å¿ƒåŠŸèƒ½æ¨¡å—",
        "â€¢ ç”¨æˆ·ç®¡ç†æ¨¡å—ï¼šæƒé™æ§åˆ¶ã€è§’è‰²ç®¡ç†",
        "â€¢ ä¸šåŠ¡å¤„ç†æ¨¡å—ï¼šæ ¸å¿ƒä¸šåŠ¡æµç¨‹",
        "â€¢ æ•°æ®åˆ†ææ¨¡å—ï¼šæŠ¥è¡¨ç»Ÿè®¡ã€è¶‹åŠ¿åˆ†æ",
        "â€¢ æŠ¥è¡¨ç”Ÿæˆæ¨¡å—ï¼šå¤šæ ¼å¼å¯¼å‡º",
        "â€¢ ç³»ç»Ÿç®¡ç†æ¨¡å—ï¼šé…ç½®ç®¡ç†ã€æ—¥å¿—å®¡è®¡",
        "",
        "æŠ€æœ¯ç‰¹ç‚¹",
        "â€¢ é«˜å¯ç”¨æ€§ï¼šé›†ç¾¤éƒ¨ç½²ã€æ•…éšœè‡ªåŠ¨åˆ‡æ¢",
        "â€¢ é«˜æ€§èƒ½ï¼šç¼“å­˜ä¼˜åŒ–ã€å¼‚æ­¥å¤„ç†",
        "â€¢ é«˜å®‰å…¨æ€§ï¼šæ•°æ®åŠ å¯†ã€æƒé™æ§åˆ¶"
    ])
    
    # ç¬¬å››éƒ¨åˆ†ï¼šé¡¹ç›®ç®¡ç†è¿‡ç¨‹ä»‹ç»
    add_section_title(prs, section_layout, "04", "é¡¹ç›®ç®¡ç†è¿‡ç¨‹ä»‹ç»")
    add_content_slide(prs, blank_layout, "å››ã€é¡¹ç›®ç®¡ç†è¿‡ç¨‹ä»‹ç»", [
        "1. ç®¡ç†æ¨¡å¼",
        "â€¢ æ•æ·å¼€å‘æ¨¡å¼ï¼ˆScrumï¼‰",
        "â€¢ ä¸¤å‘¨ä¸€ä¸ªè¿­ä»£å‘¨æœŸ",
        "â€¢ æŒç»­é›†æˆä¸æŒç»­éƒ¨ç½²ï¼ˆCI/CDï¼‰",
        "â€¢ æ¯å‘¨ä¸‰å®šæœŸä¾‹ä¼šåˆ¶åº¦",
        "â€¢ æ—¥å¸¸ç«™ä¼šï¼ˆDaily Standupï¼‰",
        "",
        "2. é¡¹ç›®ç»„æˆå‘˜",
        "â€¢ é¡¹ç›®ç»ç†ï¼š1äººï¼ˆæ•´ä½“åè°ƒï¼‰",
        "â€¢ ç³»ç»Ÿæ¶æ„å¸ˆï¼š1äººï¼ˆæŠ€æœ¯æ–¹æ¡ˆï¼‰",
        "â€¢ å‰ç«¯å·¥ç¨‹å¸ˆï¼š2-3äºº",
        "â€¢ åç«¯å·¥ç¨‹å¸ˆï¼š3-4äºº",
        "â€¢ æµ‹è¯•å·¥ç¨‹å¸ˆï¼š2äºº",
        "â€¢ UI/UXè®¾è®¡å¸ˆï¼š1-2äºº",
        "",
        "3. æ²Ÿé€šæœºåˆ¶",
        "â€¢ å‘¨ä¾‹ä¼šï¼šæ¯å‘¨ä¸‰é¡¹ç›®æ±‡æŠ¥",
        "â€¢ æ—¥å¸¸æ²Ÿé€šï¼šä¼ä¸šå¾®ä¿¡/é’‰é’‰",
        "â€¢ æ–‡æ¡£ç®¡ç†ï¼šConfluenceååŒå¹³å°",
        "â€¢ ä»£ç ç®¡ç†ï¼šGitç‰ˆæœ¬æ§åˆ¶"
    ])
    
    add_content_slide(prs, blank_layout, "å››ã€é¡¹ç›®ç®¡ç†è¿‡ç¨‹ï¼ˆç»­ï¼‰", [
        "4. è¿›åº¦è®¡åˆ’",
        "",
        "å·²å®Œæˆé˜¶æ®µï¼š",
        "â€¢ 2023å¹´4æœˆ - é¡¹ç›®å¯åŠ¨ä¸ç«‹é¡¹",
        "â€¢ 2023å¹´5æœˆ - éœ€æ±‚åˆ†æä¸è¯„å®¡",
        "â€¢ 2023å¹´6æœˆ - ç³»ç»Ÿè®¾è®¡ä¸æ¶æ„è¯„å®¡",
        "â€¢ 2023å¹´7-9æœˆ - å¼€å‘å®æ–½é˜¶æ®µ",
        "â€¢ 2023å¹´10æœˆ - ç³»ç»Ÿæµ‹è¯•ä¸bugä¿®å¤",
        "",
        "å½“å‰é˜¶æ®µï¼ˆ11æœˆï¼‰ï¼š",
        "â€¢ ç³»ç»Ÿä¼˜åŒ–ä¸æ€§èƒ½è°ƒä¼˜",
        "â€¢ ç”¨æˆ·åŸ¹è®­ä¸æ“ä½œæ‰‹å†Œ",
        "â€¢ å‡†å¤‡é¡¹ç›®éªŒæ”¶å·¥ä½œ",
        "â€¢ éªŒæ”¶æ–‡æ¡£ç¼–åˆ¶",
        "",
        "ä¸‹ä¸€é˜¶æ®µï¼ˆ12æœˆï¼‰ï¼š",
        "â€¢ æ­£å¼éªŒæ”¶ä¸è¯„å®¡",
        "â€¢ ç”Ÿäº§ç¯å¢ƒéƒ¨ç½²",
        "â€¢ é¡¹ç›®æ€»ç»“ä¸å½’æ¡£",
        "â€¢ è¿›å…¥è¿ç»´ä¿éšœæœŸ"
    ])
    
    add_content_slide(prs, blank_layout, "å››ã€é¡¹ç›®ç®¡ç†è¿‡ç¨‹ï¼ˆç»­ï¼‰", [
        "5. ç›¸å…³æˆæœç‰©",
        "",
        "éœ€æ±‚ä¸è®¾è®¡é˜¶æ®µï¼š",
        "â€¢ éœ€æ±‚è§„æ ¼è¯´æ˜ä¹¦ï¼ˆå·²å®Œæˆï¼‰",
        "â€¢ ç”¨æˆ·è°ƒç ”æŠ¥å‘Š",
        "â€¢ ç³»ç»Ÿæ¶æ„è®¾è®¡æ–‡æ¡£",
        "â€¢ æ•°æ®åº“è®¾è®¡æ–‡æ¡£",
        "â€¢ æ¥å£è®¾è®¡æ–‡æ¡£",
        "",
        "å¼€å‘ä¸æµ‹è¯•é˜¶æ®µï¼š",
        "â€¢ å®Œæ•´æºä»£ç åŠæ³¨é‡Š",
        "â€¢ å•å…ƒæµ‹è¯•æŠ¥å‘Š",
        "â€¢ é›†æˆæµ‹è¯•æŠ¥å‘Š",
        "â€¢ æ€§èƒ½æµ‹è¯•æŠ¥å‘Š",
        "â€¢ Bugä¿®å¤è®°å½•",
        "",
        "äº¤ä»˜é˜¶æ®µï¼š",
        "â€¢ ç³»ç»Ÿéƒ¨ç½²æ–‡æ¡£",
        "â€¢ ç”¨æˆ·æ“ä½œæ‰‹å†Œ",
        "â€¢ ç³»ç»Ÿè¿ç»´æ‰‹å†Œ",
        "â€¢ åŸ¹è®­ææ–™ä¸è§†é¢‘"
    ])
    
    # ç¬¬äº”éƒ¨åˆ†ï¼šäº¤ä»˜ç‰©
    add_section_title(prs, section_layout, "05", "äº¤ä»˜ç‰©")
    add_content_slide(prs, blank_layout, "äº”ã€äº¤ä»˜ç‰©", [
        "1. ç³»ç»Ÿäº¤ä»˜ç‰©",
        "â€¢ å®Œæ•´çš„ç³»ç»Ÿæºä»£ç ï¼ˆå«æ³¨é‡Šï¼‰",
        "â€¢ ç³»ç»Ÿå®‰è£…éƒ¨ç½²åŒ…",
        "â€¢ æ•°æ®åº“è„šæœ¬ï¼ˆDDL + DMLï¼‰",
        "â€¢ ç³»ç»Ÿé…ç½®æ–‡ä»¶æ¨¡æ¿",
        "â€¢ ç¬¬ä¸‰æ–¹ç»„ä»¶æ¸…å•åŠæˆæƒ",
        "",
        "2. æ–‡æ¡£äº¤ä»˜ç‰©",
        "â€¢ éœ€æ±‚è§„æ ¼è¯´æ˜ä¹¦",
        "â€¢ ç³»ç»Ÿè®¾è®¡æ–‡æ¡£ï¼ˆæ¶æ„+è¯¦ç»†è®¾è®¡ï¼‰",
        "â€¢ æ•°æ®åº“è®¾è®¡æ–‡æ¡£",
        "â€¢ APIæ¥å£æ–‡æ¡£",
        "â€¢ ç”¨æˆ·æ“ä½œæ‰‹å†Œ",
        "â€¢ ç³»ç»Ÿè¿ç»´æ‰‹å†Œ",
        "â€¢ åº”æ€¥é¢„æ¡ˆæ–‡æ¡£",
        "",
        "3. åŸ¹è®­äº¤ä»˜ç‰©",
        "â€¢ ç®¡ç†å‘˜åŸ¹è®­PPT",
        "â€¢ ç”¨æˆ·æ“ä½œåŸ¹è®­PPT",
        "â€¢ åŸ¹è®­è§†é¢‘å½•åˆ¶",
        "â€¢ å¸¸è§é—®é¢˜è§£ç­”ï¼ˆFAQï¼‰"
    ])
    
    add_content_slide(prs, blank_layout, "äº”ã€äº¤ä»˜ç‰©ï¼ˆç»­ï¼‰- æ•°æ®åˆ†ææŠ¥å‘Š", [
        "4. æ•°æ®åˆ†ææŠ¥å‘Š",
        "",
        "ç³»ç»Ÿæ€§èƒ½åˆ†æï¼š",
        "â€¢ å“åº”æ—¶é—´ï¼šå¹³å‡<500msï¼Œ95åˆ†ä½<1s",
        "â€¢ å¹¶å‘èƒ½åŠ›ï¼šæ”¯æŒ1000+å¹¶å‘ç”¨æˆ·",
        "â€¢ èµ„æºä½¿ç”¨ï¼šCPU<60%ï¼Œå†…å­˜<70%",
        "â€¢ ç³»ç»Ÿå¯ç”¨æ€§ï¼š>99.5%",
        "",
        "ä¸šåŠ¡æ•°æ®åˆ†æï¼š",
        "â€¢ ç”¨æˆ·è¡Œä¸ºåˆ†ææŠ¥å‘Š",
        "â€¢ ä¸šåŠ¡æµç¨‹æ•ˆç‡æå‡30%+",
        "â€¢ ç³»ç»Ÿæ—¥å‡è®¿é—®é‡ç»Ÿè®¡",
        "â€¢ æ ¸å¿ƒåŠŸèƒ½ä½¿ç”¨ç‡åˆ†æ",
        "",
        "è´¨é‡åˆ†ææŠ¥å‘Šï¼š",
        "â€¢ Bugç»Ÿè®¡ï¼šå·²ä¿®å¤98%ä»¥ä¸Š",
        "â€¢ ä»£ç è´¨é‡ï¼šä»£ç è¦†ç›–ç‡>80%",
        "â€¢ å®‰å…¨æ‰«æï¼šæ— é«˜å±æ¼æ´",
        "",
        "ä¼˜åŒ–å»ºè®®ï¼š",
        "â€¢ åç»­åŠŸèƒ½æ‰©å±•è§„åˆ’",
        "â€¢ æ€§èƒ½æŒç»­ä¼˜åŒ–æ–¹æ¡ˆ",
        "â€¢ ç”¨æˆ·ä½“éªŒæ”¹è¿›å»ºè®®"
    ])
    
    # ç¬¬å…­éƒ¨åˆ†ï¼šé¡¹ç›®éªŒæ”¶
    add_section_title(prs, section_layout, "06", "é¡¹ç›®éªŒæ”¶")
    add_content_slide(prs, blank_layout, "å…­ã€é¡¹ç›®éªŒæ”¶", [
        "éªŒæ”¶å‡†å¤‡å·¥ä½œ",
        "",
        "1. éªŒæ”¶æ–‡æ¡£å‡†å¤‡",
        "â€¢ é¡¹ç›®éªŒæ”¶ç”³è¯·æŠ¥å‘Š",
        "â€¢ é¡¹ç›®æ€»ç»“æŠ¥å‘Š",
        "â€¢ ç³»ç»Ÿæµ‹è¯•æŠ¥å‘Šæ±‡æ€»",
        "â€¢ ç”¨æˆ·éªŒæ”¶æµ‹è¯•ï¼ˆUATï¼‰æŠ¥å‘Š",
        "â€¢ é¡¹ç›®å˜æ›´è®°å½•æ¸…å•",
        "â€¢ é—®é¢˜è·Ÿè¸ªä¸è§£å†³è®°å½•",
        "",
        "2. éªŒæ”¶ç¯å¢ƒå‡†å¤‡",
        "â€¢ éªŒæ”¶æµ‹è¯•ç¯å¢ƒæ­å»ºå®Œæˆ",
        "â€¢ éªŒæ”¶æ•°æ®å‡†å¤‡ä¸è„±æ•",
        "â€¢ éªŒæ”¶æ¼”ç¤ºè„šæœ¬ç¼–å†™",
        "â€¢ éªŒæ”¶å›¢é˜Ÿäººå‘˜ç»„ç»‡",
        "",
        "3. éªŒæ”¶æ ‡å‡†",
        "â€¢ åŠŸèƒ½å®Œæ•´æ€§è¾¾æ ‡ï¼ˆ100%ï¼‰",
        "â€¢ æ€§èƒ½æŒ‡æ ‡è¾¾æ ‡ï¼ˆæ»¡è¶³SLAè¦æ±‚ï¼‰",
        "â€¢ å®‰å…¨æ€§è¯„ä¼°é€šè¿‡",
        "â€¢ æ–‡æ¡£å®Œæ•´æ€§è¾¾æ ‡"
    ])
    
    add_content_slide(prs, blank_layout, "å…­ã€é¡¹ç›®éªŒæ”¶ï¼ˆç»­ï¼‰", [
        "éªŒæ”¶å†…å®¹",
        "",
        "1. åŠŸèƒ½éªŒæ”¶",
        "â€¢ æ ¸å¿ƒä¸šåŠ¡åŠŸèƒ½éªŒè¯",
        "â€¢ ç”¨æˆ·æƒé™ç®¡ç†éªŒè¯",
        "â€¢ æ•°æ®å¤„ç†å‡†ç¡®æ€§éªŒè¯",
        "â€¢ å¼‚å¸¸å¤„ç†æœºåˆ¶éªŒè¯",
        "â€¢ ç³»ç»Ÿé›†æˆæ¥å£éªŒè¯",
        "",
        "2. æ€§èƒ½éªŒæ”¶",
        "â€¢ å“åº”æ—¶é—´æµ‹è¯•",
        "â€¢ å¹¶å‘å‹åŠ›æµ‹è¯•",
        "â€¢ å¤§æ•°æ®é‡å¤„ç†æµ‹è¯•",
        "â€¢ é•¿æ—¶é—´ç¨³å®šæ€§æµ‹è¯•",
        "",
        "3. å®‰å…¨æ€§éªŒæ”¶",
        "â€¢ èº«ä»½è®¤è¯ä¸æˆæƒ",
        "â€¢ æ•°æ®åŠ å¯†ä¼ è¾“",
        "â€¢ SQLæ³¨å…¥é˜²æŠ¤",
        "â€¢ XSSæ”»å‡»é˜²æŠ¤",
        "â€¢ æ•æ„Ÿä¿¡æ¯ä¿æŠ¤",
        "",
        "4. æ˜“ç”¨æ€§éªŒæ”¶",
        "â€¢ ç”¨æˆ·ç•Œé¢å‹å¥½æ€§",
        "â€¢ æ“ä½œæµç•…åº¦",
        "â€¢ å¸®åŠ©æ–‡æ¡£å®Œæ•´æ€§"
    ])
    
    add_content_slide(prs, blank_layout, "å…­ã€é¡¹ç›®éªŒæ”¶ï¼ˆç»­ï¼‰", [
        "éªŒæ”¶è®¡åˆ’ä¸æ—¶é—´å®‰æ’",
        "",
        "ç¬¬ä¸€é˜¶æ®µï¼šå†…éƒ¨éªŒæ”¶ï¼ˆç¬¬1å‘¨ï¼‰",
        "â€¢ å¼€å‘å›¢é˜Ÿè‡ªæµ‹å®Œæˆ",
        "â€¢ æµ‹è¯•å›¢é˜Ÿå…¨é¢æµ‹è¯•",
        "â€¢ é—®é¢˜ä¿®å¤ä¸å›å½’æµ‹è¯•",
        "",
        "ç¬¬äºŒé˜¶æ®µï¼šç”¨æˆ·éªŒæ”¶ï¼ˆç¬¬2-3å‘¨ï¼‰",
        "â€¢ ç”¨æˆ·ä»£è¡¨å‚ä¸UATæµ‹è¯•",
        "â€¢ æ”¶é›†ç”¨æˆ·åé¦ˆæ„è§",
        "â€¢ ä¼˜åŒ–è°ƒæ•´ä¸äºŒæ¬¡éªŒè¯",
        "",
        "ç¬¬ä¸‰é˜¶æ®µï¼šæ­£å¼éªŒæ”¶ï¼ˆç¬¬4å‘¨ï¼‰",
        "â€¢ æäº¤éªŒæ”¶ç”³è¯·",
        "â€¢ ç»„ç»‡éªŒæ”¶è¯„å®¡ä¼š",
        "â€¢ ç°åœºæ¼”ç¤ºä¸ç­”è¾©",
        "â€¢ éªŒæ”¶ä¸“å®¶ç»„è¯„å®¡",
        "â€¢ ç­¾ç½²éªŒæ”¶æŠ¥å‘Š",
        "",
        "éªŒæ”¶é€šè¿‡æ ‡å‡†ï¼š",
        "â€¢ æ‰€æœ‰åŠŸèƒ½æ­£å¸¸è¿è¡Œ",
        "â€¢ æ€§èƒ½æŒ‡æ ‡è¾¾åˆ°è¦æ±‚",
        "â€¢ æ— é—ç•™é«˜å±é—®é¢˜",
        "â€¢ æ–‡æ¡£èµ„æ–™å®Œæ•´é½å…¨",
        "â€¢ éªŒæ”¶ä¸“å®¶ç»„ä¸€è‡´é€šè¿‡"
    ])
    
    add_content_slide(prs, blank_layout, "å…­ã€é¡¹ç›®éªŒæ”¶ï¼ˆç»­ï¼‰", [
        "éªŒæ”¶é£é™©ä¸åº”å¯¹",
        "",
        "æ½œåœ¨é£é™©ï¼š",
        "â€¢ åŠŸèƒ½ç¼ºé™·ï¼šå°æ¦‚ç‡é—ç•™bug",
        "   åº”å¯¹ï¼šå»ºç«‹å¿«é€Ÿå“åº”æœºåˆ¶",
        "",
        "â€¢ æ€§èƒ½ä¸è¾¾æ ‡ï¼šé«˜å¹¶å‘åœºæ™¯å‹åŠ›",
        "   åº”å¯¹ï¼šæå‰ä¼˜åŒ–ï¼Œé¢„ç•™æ€§èƒ½ä½™é‡",
        "",
        "â€¢ ç”¨æˆ·åŸ¹è®­ä¸è¶³ï¼šæ“ä½œä¸ç†Ÿç»ƒ",
        "   åº”å¯¹ï¼šåŠ å¼ºåŸ¹è®­ï¼Œæä¾›ç°åœºæ”¯æŒ",
        "",
        "â€¢ æ–‡æ¡£ä¸å®Œæ•´ï¼šç¼ºå°‘éƒ¨åˆ†æ–‡æ¡£",
        "   åº”å¯¹ï¼šæ–‡æ¡£æ£€æŸ¥æ¸…å•ï¼Œé€é¡¹ç¡®è®¤",
        "",
        "éªŒæ”¶æˆåŠŸä¿éšœæªæ–½ï¼š",
        "â€¢ å»ºç«‹éªŒæ”¶å€’è®¡æ—¶æœºåˆ¶",
        "â€¢ æ¯æ—¥éªŒæ”¶å‡†å¤‡è¿›åº¦åŒæ­¥",
        "â€¢ ä¸“äººè´Ÿè´£éªŒæ”¶åè°ƒå·¥ä½œ",
        "â€¢ é¢„ç•™å……è¶³çš„é—®é¢˜ä¿®å¤æ—¶é—´",
        "â€¢ å‡†å¤‡è¯¦ç»†çš„æ¼”ç¤ºæ–¹æ¡ˆ",
        "â€¢ æå‰ä¸éªŒæ”¶æ–¹å……åˆ†æ²Ÿé€š"
    ])
    
    # æ€»ç»“é¡µ
    add_content_slide(prs, blank_layout, "é¡¹ç›®æ€»ç»“ä¸å±•æœ›", [
        "é¡¹ç›®æˆæœ",
        "â€¢ æŒ‰æœŸå®Œæˆç³»ç»Ÿå¼€å‘ä¸æµ‹è¯•",
        "â€¢ ç³»ç»ŸåŠŸèƒ½å®Œæ•´ï¼Œæ€§èƒ½ç¨³å®š",
        "â€¢ å›¢é˜Ÿåä½œé¡ºç•…ï¼Œç®¡ç†è§„èŒƒ",
        "â€¢ æ–‡æ¡£èµ„æ–™å®Œæ•´ï¼Œäº¤ä»˜ç‰©é½å…¨",
        "",
        "é¡¹ç›®äº®ç‚¹",
        "â€¢ é‡‡ç”¨å…ˆè¿›æŠ€æœ¯æ¶æ„ï¼Œæ˜“äºæ‰©å±•",
        "â€¢ ç”¨æˆ·ç•Œé¢å‹å¥½ï¼Œæ“ä½œä¾¿æ·",
        "â€¢ ç³»ç»Ÿæ€§èƒ½ä¼˜å¼‚ï¼Œå“åº”å¿«é€Ÿ",
        "â€¢ å®‰å…¨æœºåˆ¶å®Œå–„ï¼Œæ•°æ®å¯é ",
        "",
        "åç»­è®¡åˆ’",
        "â€¢ é¡ºåˆ©é€šè¿‡é¡¹ç›®éªŒæ”¶",
        "â€¢ ç¨³å®šä¸Šçº¿è¿è¡Œ",
        "â€¢ æŒç»­ä¼˜åŒ–æ”¹è¿›",
        "â€¢ ç”¨æˆ·åŸ¹è®­ä¸æ”¯æŒ",
        "â€¢ äºŒæœŸåŠŸèƒ½è§„åˆ’",
        "",
        "è‡´è°¢",
        "â€¢ æ„Ÿè°¢é¡¹ç›®ç»„å…¨ä½“æˆå‘˜çš„è¾›å‹¤ä»˜å‡º",
        "â€¢ æ„Ÿè°¢ç”¨æˆ·æ–¹çš„å¤§åŠ›æ”¯æŒä¸é…åˆ",
        "â€¢ æ„Ÿè°¢å„çº§é¢†å¯¼çš„å…³å¿ƒä¸æŒ‡å¯¼"
    ])
    
    # ç»“æŸé¡µ
    slide = prs.slides.add_slide(blank_layout)
    end_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(6), Inches(1.5))
    end_frame = end_box.text_frame
    end_frame.text = "è°¢è°¢ï¼\n\nQ & A"
    for para in end_frame.paragraphs:
        para.font.name = CHINESE_FONT
        para.font.size = Pt(48)
        para.font.bold = True
        para.font.color.rgb = MAIN_COLOR
        para.alignment = PP_ALIGN.CENTER
        set_font(para, CHINESE_FONT)
    
    # ä¿å­˜æ¼”ç¤ºæ–‡ç¨¿
    output_file = '/workspace/Wednesday_Report_New.pptx'
    prs.save(output_file)
    print("âœ… PPTç”ŸæˆæˆåŠŸï¼")
    print("ğŸ“„ æ–‡ä»¶è·¯å¾„ï¼š{}".format(output_file))
    print("ğŸ¨ ä½¿ç”¨äº†æµ·æ˜Ÿè‚²æ•°å­—åŒ–ç³»ç»Ÿæ¨¡æ¿æ ·å¼")
    return output_file

def add_custom_title_slide(prs, slide, title, subtitle, date):
    """æ‰‹åŠ¨åœ¨æ ‡é¢˜é¡µæ·»åŠ æ–‡æœ¬"""
    # æ ‡é¢˜
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.name = CHINESE_FONT
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = MAIN_COLOR
    title_para.alignment = PP_ALIGN.CENTER
    set_font(title_para, CHINESE_FONT)
    
    # å‰¯æ ‡é¢˜
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.name = CHINESE_FONT
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.color.rgb = RGBColor(89, 89, 89)
    subtitle_para.alignment = PP_ALIGN.CENTER
    set_font(subtitle_para, CHINESE_FONT)
    
    # æ—¥æœŸ
    date_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.5))
    date_frame = date_box.text_frame
    date_frame.text = date
    date_para = date_frame.paragraphs[0]
    date_para.font.name = CHINESE_FONT
    date_para.font.size = Pt(18)
    date_para.font.color.rgb = RGBColor(89, 89, 89)
    date_para.alignment = PP_ALIGN.CENTER
    set_font(date_para, CHINESE_FONT)

def add_section_title(prs, layout, number, title):
    """æ·»åŠ èŠ‚æ ‡é¢˜"""
    slide = prs.slides.add_slide(layout)
    
    # æ¸…ç©ºå ä½ç¬¦
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            for para in shape.text_frame.paragraphs:
                if len(para.runs) > 0:
                    para.clear()
    
    # æ·»åŠ æ•°å­—
    num_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1.2))
    num_frame = num_box.text_frame
    num_frame.text = number
    num_para = num_frame.paragraphs[0]
    num_para.font.name = CHINESE_FONT
    num_para.font.size = Pt(88)
    num_para.font.bold = True
    num_para.font.color.rgb = MAIN_COLOR
    num_para.alignment = PP_ALIGN.CENTER
    set_font(num_para, CHINESE_FONT)
    
    # æ·»åŠ æ ‡é¢˜
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.name = CHINESE_FONT
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = MAIN_COLOR
    title_para.alignment = PP_ALIGN.CENTER
    set_font(title_para, CHINESE_FONT)

def add_content_slide(prs, layout, title, content_list):
    """æ·»åŠ å†…å®¹é¡µ"""
    slide = prs.slides.add_slide(layout)
    
    # æ·»åŠ æ ‡é¢˜
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.name = CHINESE_FONT
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = MAIN_COLOR
    set_font(title_para, CHINESE_FONT)
    
    # æ·»åŠ å†…å®¹
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(4.8))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    for i, line in enumerate(content_list):
        if i == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()
        
        p.text = line
        p.font.name = CHINESE_FONT
        p.font.size = Pt(14)
        p.space_after = Pt(4)
        
        if line.startswith('â€¢'):
            p.level = 0
            p.font.size = Pt(14)
        elif line.startswith('   â€¢'):
            p.level = 1
            p.font.size = Pt(13)
        elif line.startswith('   '):
            p.level = 1
            p.font.size = Pt(13)
        
        # åŠ ç²—æ ‡é¢˜è¡Œ
        if line and not line.startswith(' ') and not line.startswith('â€¢') and not line.startswith('å·²') and not line.startswith('å½“') and not line.startswith('ä¸‹') and not line.startswith('ç¬¬') and ':' not in line and 'ï¼š' not in line and line != "" and not line.startswith('éªŒæ”¶'):
            p.font.bold = True
            p.font.size = Pt(16)
            p.font.color.rgb = MAIN_COLOR
        
        set_font(p, CHINESE_FONT)

if __name__ == "__main__":
    create_presentation_from_template()
