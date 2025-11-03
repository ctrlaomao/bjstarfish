#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement

CHINESE_FONT = "WenQuanYi Zen Hei"

def set_font(paragraph, font_name):
    """设置段落字体，包括东亚字体"""
    try:
        for run in paragraph.runs:
            run.font.name = font_name
            if hasattr(run.font, '_element'):
                rPr = run.font._element.get_or_add_rPr()
                ea = OxmlElement('a:ea')
                ea.set('typeface', font_name)
                rPr.append(ea)
                cs = OxmlElement('a:cs')
                cs.set('typeface', font_name)
                rPr.append(cs)
    except:
        pass

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 标题页
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "周三项目汇报"
    title_para = title_frame.paragraphs[0]
    title_para.font.name = CHINESE_FONT
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    title_para.alignment = PP_ALIGN.CENTER
    set_font(title_para, CHINESE_FONT)
    
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "胜达讯项目团队"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.name = CHINESE_FONT
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = RGBColor(89, 89, 89)
    subtitle_para.alignment = PP_ALIGN.CENTER
    set_font(subtitle_para, CHINESE_FONT)
    
    # 第一部分：项目规划
    add_section_slide(prs, "一、项目规划", [
        "项目目标",
        "• 构建高效、可靠的业务系统",
        "• 提升业务处理效率和用户体验",
        "• 实现数据驱动的决策支持",
        "• 保障系统安全性和稳定性",
        "",
        "项目范围",
        "• 系统功能模块开发",
        "• 数据库设计与实现",
        "• 接口开发与集成",
        "• 系统测试与部署",
        "• 用户培训与文档编制",
        "",
        "项目周期",
        "• 启动时间：2023年4月",
        "• 开发周期：2023年4月 - 10月",
        "• 当前状态：准备验收阶段"
    ])
    
    # 第二部分：业务流程
    add_section_slide(prs, "二、业务流程", [
        "核心业务流程",
        "",
        "1. 需求收集与分析",
        "   • 用户需求调研",
        "   • 业务流程梳理",
        "   • 需求文档编写",
        "   • 需求评审与确认",
        "",
        "2. 系统设计",
        "   • 总体架构设计",
        "   • 数据库详细设计",
        "   • 接口设计与规范",
        "   • 安全机制设计",
        "",
        "3. 开发实现",
        "   • 模块化开发",
        "   • 代码审查",
        "   • 单元测试",
        "   • 集成测试"
    ])
    
    add_section_slide(prs, "二、业务流程（续）", [
        "4. 系统测试",
        "   • 功能测试",
        "   • 性能测试",
        "   • 安全测试",
        "   • 用户验收测试（UAT）",
        "",
        "5. 部署上线",
        "   • 生产环境准备",
        "   • 系统部署实施",
        "   • 数据迁移与验证",
        "   • 切换方案执行",
        "",
        "6. 运维支持",
        "   • 7×24小时系统监控",
        "   • 快速问题响应机制",
        "   • 持续优化与改进",
        "   • 定期巡检与维护"
    ])
    
    # 第三部分：系统演示
    add_section_slide(prs, "三、系统演示", [
        "系统架构",
        "• 前端展示层：响应式Web界面",
        "• 业务逻辑层：RESTful API服务",
        "• 数据访问层：ORM框架封装",
        "• 数据存储层：关系型数据库 + 缓存",
        "• 基础设施层：云服务支持",
        "",
        "核心功能模块",
        "• 用户管理模块：权限控制、角色管理",
        "• 业务处理模块：核心业务流程",
        "• 数据分析模块：报表统计、趋势分析",
        "• 报表生成模块：多格式导出",
        "• 系统管理模块：配置管理、日志审计",
        "",
        "技术特点",
        "• 高可用性：集群部署、故障自动切换",
        "• 高性能：缓存优化、异步处理",
        "• 高安全性：数据加密、权限控制"
    ])
    
    # 第四部分：项目管理过程介绍
    add_section_slide(prs, "四、项目管理过程介绍", [
        "1. 管理模式",
        "• 敏捷开发模式（Scrum）",
        "• 两周一个迭代周期",
        "• 持续集成与持续部署（CI/CD）",
        "• 每周三定期例会制度",
        "• 日常站会（Daily Standup）",
        "",
        "2. 项目组成员",
        "• 项目经理：1人（整体协调）",
        "• 系统架构师：1人（技术方案）",
        "• 前端工程师：2-3人",
        "• 后端工程师：3-4人",
        "• 测试工程师：2人",
        "• UI/UX设计师：1-2人",
        "",
        "3. 沟通机制",
        "• 周例会：每周三项目汇报",
        "• 日常沟通：企业微信/钉钉",
        "• 文档管理：Confluence协同平台",
        "• 代码管理：Git版本控制"
    ])
    
    add_section_slide(prs, "四、项目管理过程（续）", [
        "4. 进度计划",
        "",
        "已完成阶段：",
        "• 2023年4月 - 项目启动与立项",
        "• 2023年5月 - 需求分析与评审",
        "• 2023年6月 - 系统设计与架构评审",
        "• 2023年7-9月 - 开发实施阶段",
        "• 2023年10月 - 系统测试与bug修复",
        "",
        "当前阶段（11月）：",
        "• 系统优化与性能调优",
        "• 用户培训与操作手册",
        "• 准备项目验收工作",
        "• 验收文档编制",
        "",
        "下一阶段（12月）：",
        "• 正式验收与评审",
        "• 生产环境部署",
        "• 项目总结与归档",
        "• 进入运维保障期"
    ])
    
    add_section_slide(prs, "四、项目管理过程（续）", [
        "5. 相关成果物",
        "",
        "需求与设计阶段：",
        "• 需求规格说明书（已完成）",
        "• 用户调研报告",
        "• 系统架构设计文档",
        "• 数据库设计文档",
        "• 接口设计文档",
        "",
        "开发与测试阶段：",
        "• 完整源代码及注释",
        "• 单元测试报告",
        "• 集成测试报告",
        "• 性能测试报告",
        "• Bug修复记录",
        "",
        "交付阶段：",
        "• 系统部署文档",
        "• 用户操作手册",
        "• 系统运维手册",
        "• 培训材料与视频"
    ])
    
    # 第五部分：交付物
    add_section_slide(prs, "五、交付物", [
        "1. 系统交付物",
        "• 完整的系统源代码（含注释）",
        "• 系统安装部署包",
        "• 数据库脚本（DDL + DML）",
        "• 系统配置文件模板",
        "• 第三方组件清单及授权",
        "",
        "2. 文档交付物",
        "• 需求规格说明书",
        "• 系统设计文档（架构+详细设计）",
        "• 数据库设计文档",
        "• API接口文档",
        "• 用户操作手册",
        "• 系统运维手册",
        "• 应急预案文档",
        "",
        "3. 培训交付物",
        "• 管理员培训PPT",
        "• 用户操作培训PPT",
        "• 培训视频录制",
        "• 常见问题解答（FAQ）"
    ])
    
    add_section_slide(prs, "五、交付物（续）- 数据分析报告", [
        "4. 数据分析报告",
        "",
        "系统性能分析：",
        "• 响应时间：平均<500ms，95分位<1s",
        "• 并发能力：支持1000+并发用户",
        "• 资源使用：CPU<60%，内存<70%",
        "• 系统可用性：>99.5%",
        "",
        "业务数据分析：",
        "• 用户行为分析报告",
        "• 业务流程效率提升30%+",
        "• 系统日均访问量统计",
        "• 核心功能使用率分析",
        "",
        "质量分析报告：",
        "• Bug统计：已修复98%以上",
        "• 代码质量：代码覆盖率>80%",
        "• 安全扫描：无高危漏洞",
        "",
        "优化建议：",
        "• 后续功能扩展规划",
        "• 性能持续优化方案",
        "• 用户体验改进建议"
    ])
    
    # 第六部分：项目验收
    add_section_slide(prs, "六、项目验收", [
        "验收准备工作",
        "",
        "1. 验收文档准备",
        "• 项目验收申请报告",
        "• 项目总结报告",
        "• 系统测试报告汇总",
        "• 用户验收测试（UAT）报告",
        "• 项目变更记录清单",
        "• 问题跟踪与解决记录",
        "",
        "2. 验收环境准备",
        "• 验收测试环境搭建完成",
        "• 验收数据准备与脱敏",
        "• 验收演示脚本编写",
        "• 验收团队人员组织",
        "",
        "3. 验收标准",
        "• 功能完整性达标（100%）",
        "• 性能指标达标（满足SLA要求）",
        "• 安全性评估通过",
        "• 文档完整性达标"
    ])
    
    add_section_slide(prs, "六、项目验收（续）", [
        "验收内容",
        "",
        "1. 功能验收",
        "• 核心业务功能验证",
        "• 用户权限管理验证",
        "• 数据处理准确性验证",
        "• 异常处理机制验证",
        "• 系统集成接口验证",
        "",
        "2. 性能验收",
        "• 响应时间测试",
        "• 并发压力测试",
        "• 大数据量处理测试",
        "• 长时间稳定性测试",
        "",
        "3. 安全性验收",
        "• 身份认证与授权",
        "• 数据加密传输",
        "• SQL注入防护",
        "• XSS攻击防护",
        "• 敏感信息保护",
        "",
        "4. 易用性验收",
        "• 用户界面友好性",
        "• 操作流畅度",
        "• 帮助文档完整性"
    ])
    
    add_section_slide(prs, "六、项目验收（续）", [
        "验收计划与时间安排",
        "",
        "第一阶段：内部验收（第1周）",
        "• 开发团队自测完成",
        "• 测试团队全面测试",
        "• 问题修复与回归测试",
        "",
        "第二阶段：用户验收（第2-3周）",
        "• 用户代表参与UAT测试",
        "• 收集用户反馈意见",
        "• 优化调整与二次验证",
        "",
        "第三阶段：正式验收（第4周）",
        "• 提交验收申请",
        "• 组织验收评审会",
        "• 现场演示与答辩",
        "• 验收专家组评审",
        "• 签署验收报告",
        "",
        "验收通过标准：",
        "• 所有功能正常运行",
        "• 性能指标达到要求",
        "• 无遗留高危问题",
        "• 文档资料完整齐全",
        "• 验收专家组一致通过"
    ])
    
    add_section_slide(prs, "六、项目验收（续）", [
        "验收风险与应对",
        "",
        "潜在风险：",
        "• 功能缺陷：小概率遗留bug",
        "   应对：建立快速响应机制",
        "",
        "• 性能不达标：高并发场景压力",
        "   应对：提前优化，预留性能余量",
        "",
        "• 用户培训不足：操作不熟练",
        "   应对：加强培训，提供现场支持",
        "",
        "• 文档不完整：缺少部分文档",
        "   应对：文档检查清单，逐项确认",
        "",
        "验收成功保障措施：",
        "• 建立验收倒计时机制",
        "• 每日验收准备进度同步",
        "• 专人负责验收协调工作",
        "• 预留充足的问题修复时间",
        "• 准备详细的演示方案",
        "• 提前与验收方充分沟通"
    ])
    
    # 总结页
    add_section_slide(prs, "项目总结与展望", [
        "项目成果",
        "• 按期完成系统开发与测试",
        "• 系统功能完整，性能稳定",
        "• 团队协作顺畅，管理规范",
        "• 文档资料完整，交付物齐全",
        "",
        "项目亮点",
        "• 采用先进技术架构，易于扩展",
        "• 用户界面友好，操作便捷",
        "• 系统性能优异，响应快速",
        "• 安全机制完善，数据可靠",
        "",
        "后续计划",
        "• 顺利通过项目验收",
        "• 稳定上线运行",
        "• 持续优化改进",
        "• 用户培训与支持",
        "• 二期功能规划",
        "",
        "致谢",
        "• 感谢项目组全体成员的辛勤付出",
        "• 感谢用户方的大力支持与配合",
        "• 感谢各级领导的关心与指导"
    ])
    
    # 结束页
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    end_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
    end_frame = end_box.text_frame
    end_frame.text = "谢谢！\n\nQ & A"
    for para in end_frame.paragraphs:
        para.font.name = CHINESE_FONT
        para.font.size = Pt(44)
        para.font.bold = True
        para.font.color.rgb = RGBColor(0, 51, 102)
        para.alignment = PP_ALIGN.CENTER
        set_font(para, CHINESE_FONT)
    
    # 保存演示文稿
    output_file = '/workspace/Wednesday_Report.pptx'
    prs.save(output_file)
    print("PPT生成成功！")
    print("文件路径：{}".format(output_file))
    print("字体设置：{}".format(CHINESE_FONT))
    return output_file

def add_section_slide(prs, title, content_list):
    """添加内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.name = CHINESE_FONT
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    set_font(title_para, CHINESE_FONT)
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    for i, line in enumerate(content_list):
        if i == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()
        
        p.text = line
        p.font.name = CHINESE_FONT
        p.font.size = Pt(15)
        p.space_after = Pt(5)
        
        if line.startswith('•'):
            p.level = 0
            p.font.size = Pt(15)
        elif line.startswith('   •'):
            p.level = 1
            p.font.size = Pt(13)
        elif line.startswith('   '):
            p.level = 1
            p.font.size = Pt(13)
        
        if line and not line.startswith(' ') and not line.startswith('•') and not line.startswith('已') and not line.startswith('当') and not line.startswith('下') and not line.startswith('第') and ':' not in line and '：' not in line and line != "":
            p.font.bold = True
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(0, 51, 102)
        
        set_font(p, CHINESE_FONT)

if __name__ == "__main__":
    create_presentation()
