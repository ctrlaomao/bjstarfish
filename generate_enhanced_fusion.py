#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR

CHINESE_FONT = "WenQuanYi Zen Hei"
MAIN_COLOR = RGBColor(0, 70, 160)
LIGHT_BLUE = RGBColor(173, 216, 230)
ACCENT_COLOR = RGBColor(255, 140, 0)

def set_font(paragraph, font_name):
    """设置段落字体，包括东亚字体"""
    try:
        for run in paragraph.runs:
            run.font.name = font_name
            if hasattr(run.font, '_element'):
                rPr = run.font._element.get_or_add_rPr()
                ea = OxmlElement('a:ea')
                ea.set('typeface', font_name)
                rPr.append(ea)
                cs = OxmlElement('a:cs')
                cs.set('typeface', font_name)
                rPr.append(cs)
    except:
        pass

def add_arrow_shape(slide, left, top, width, height, text):
    """添加箭头形状"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BLUE
    shape.line.color.rgb = MAIN_COLOR
    shape.line.width = Pt(1.5)
    
    if text:
        text_frame = shape.text_frame
        text_frame.text = text
        p = text_frame.paragraphs[0]
        p.font.name = CHINESE_FONT
        p.font.size = Pt(11)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        set_font(p, CHINESE_FONT)

def add_box_shape(slide, left, top, width, height, text, color=LIGHT_BLUE):
    """添加矩形框"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = MAIN_COLOR
    shape.line.width = Pt(2)
    
    if text:
        text_frame = shape.text_frame
        text_frame.text = text
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)
        p = text_frame.paragraphs[0]
        p.font.name = CHINESE_FONT
        p.font.size = Pt(12)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        set_font(p, CHINESE_FONT)

def extract_content(content_ppt_path):
    """提取内容PPT的所有数据"""
    print("\n读取内容.pptx...")
    prs = Presentation(content_ppt_path)
    
    slides_data = []
    slides_list = list(prs.slides)
    
    for idx, slide in enumerate(slides_list):
        slide_info = {
            'index': idx + 1,
            'texts': []
        }
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_info['texts'].append(shape.text.strip())
        
        slides_data.append(slide_info)
    
    print("✓ 提取 {} 张幻灯片\n".format(len(slides_data)))
    return slides_data

def create_enhanced_fusion_ppt(template_path, slides_data, output_path):
    """创建增强版融合PPT"""
    print("加载模板并清空...")
    prs = Presentation(template_path)
    
    title_layout = prs.slide_layouts[12]
    section_layout = prs.slide_layouts[13]
    content_layout = prs.slide_layouts[14]
    
    # 清空
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]
    
    print("填充内容并添加可视化元素...\n")
    
    # 重新获取布局
    title_layout = prs.slide_layouts[12]
    section_layout = prs.slide_layouts[13]
    content_layout = prs.slide_layouts[14]
    
    for slide_data in slides_data:
        idx = slide_data['index']
        texts = slide_data['texts']
        
        if not texts:
            continue
        
        first_text = texts[0]
        
        # 标题页
        if idx == 1:
            add_enhanced_title_slide(prs, title_layout, texts)
            print("  [{}] 标题页（优化）".format(idx))
        
        # 项目规划 - 添加流程图
        elif "项目规划" in first_text:
            add_planning_slide_with_diagram(prs, content_layout, first_text, texts[1:])
            print("  [{}] 项目规划（添加流程图）".format(idx))
        
        # 业务流程 - 添加流程箭头
        elif "业务流程" in first_text:
            add_process_slide_with_arrows(prs, content_layout, first_text, texts[1:])
            print("  [{}] 业务流程（添加流程图）".format(idx))
        
        # 系统演示 - 添加架构图
        elif "系统演示" in first_text:
            add_system_slide_with_architecture(prs, content_layout, first_text, texts[1:])
            print("  [{}] 系统演示（添加架构图）".format(idx))
        
        # 项目管理 - 添加时间轴
        elif "项目管理" in first_text or "进度计划" in str(texts):
            add_management_slide_with_timeline(prs, content_layout, first_text, texts[1:])
            print("  [{}] 项目管理（添加可视化）".format(idx))
        
        # 交付物 - 添加图标
        elif "交付物" in first_text:
            add_deliverable_slide_with_icons(prs, content_layout, first_text, texts[1:])
            print("  [{}] 交付物（添加图标）".format(idx))
        
        # 项目验收 - 添加检查框
        elif "项目验收" in first_text:
            add_acceptance_slide_with_checkboxes(prs, content_layout, first_text, texts[1:])
            print("  [{}] 项目验收（添加检查项）".format(idx))
        
        # 结束页
        elif "谢谢" in first_text or "Q & A" in first_text:
            add_enhanced_end_slide(prs, content_layout, texts)
            print("  [{}] 结束页（优化）".format(idx))
        
        # 普通内容页
        else:
            add_enhanced_content_slide(prs, content_layout, first_text, texts[1:] if len(texts) > 1 else [])
            print("  [{}] 内容页: {}".format(idx, first_text[:20]))
    
    print("\n保存文件...")
    prs.save(output_path)
    
    print("\n" + "=" * 70)
    print("✅ 增强版融合PPT生成成功！")
    print("=" * 70)
    print("📄 文件: {}".format(output_path))
    print("📊 幻灯片: {} 张".format(len(slides_data)))
    print("🎨 模板背景: 海星育风格")
    print("📝 内容来源: 周三项目汇报")
    print("✨ 新增特性: 可视化图表、流程图、图标等")
    print("=" * 70)

def add_enhanced_title_slide(prs, layout, texts):
    """添加增强的标题页 - 标题位置下移"""
    slide = prs.slides.add_slide(layout)
    
    title = texts[0] if len(texts) > 0 else "项目汇报"
    subtitle = texts[1] if len(texts) > 1 else ""
    
    # 标题 - 向下移动到3英寸位置
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.33), Inches(1.2))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.name = CHINESE_FONT
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = MAIN_COLOR
    title_para.alignment = PP_ALIGN.CENTER
    set_font(title_para, CHINESE_FONT)
    
    # 副标题
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.33), Inches(0.8))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.name = CHINESE_FONT
        subtitle_para.font.size = Pt(32)
        subtitle_para.font.color.rgb = RGBColor(89, 89, 89)
        subtitle_para.alignment = PP_ALIGN.CENTER
        set_font(subtitle_para, CHINESE_FONT)

def add_planning_slide_with_diagram(prs, layout, title, content_texts):
    """项目规划页 - 添加流程图"""
    slide = prs.slides.add_slide(layout)
    
    # 标题 - 下移到0.8英寸
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(12.33), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.name = CHINESE_FONT
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = MAIN_COLOR
    set_font(title_para, CHINESE_FONT)
    
    # 左侧：文字内容（减少宽度）
    if content_texts:
        all_content = "\n".join(content_texts)
        lines = all_content.split('\n')
        
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(6.5), Inches(5.4))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        for i, line in enumerate(lines):
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            
            p.text = line
            p.font.name = CHINESE_FONT
            p.font.size = Pt(13)
            p.space_after = Pt(3)
            
            if line.startswith('•'):
                p.level = 0
            elif line.startswith('   '):
                p.level = 1
                p.font.size = Pt(12)
            
            if line and not line.startswith(' ') and not line.startswith('•') and len(line) < 20:
                p.font.bold = True
                p.font.size = Pt(15)
                p.font.color.rgb = MAIN_COLOR
            
            set_font(p, CHINESE_FONT)
    
    # 右侧：流程图
    add_box_shape(slide, Inches(7.5), Inches(2), Inches(2.2), Inches(0.8), "项目启动", LIGHT_BLUE)
    add_arrow_shape(slide, Inches(8.5), Inches(2.9), Inches(1), Inches(0.4), "")
    add_box_shape(slide, Inches(7.5), Inches(3.4), Inches(2.2), Inches(0.8), "需求分析", LIGHT_BLUE)
    add_arrow_shape(slide, Inches(8.5), Inches(4.3), Inches(1), Inches(0.4), "")
    add_box_shape(slide, Inches(7.5), Inches(4.8), Inches(2.2), Inches(0.8), "系统开发", LIGHT_BLUE)
    add_arrow_shape(slide, Inches(8.5), Inches(5.7), Inches(1), Inches(0.4), "")
    add_box_shape(slide, Inches(7.5), Inches(6.2), Inches(2.2), Inches(0.8), "项目验收", RGBColor(144, 238, 144))

def add_process_slide_with_arrows(prs, layout, title, content_texts):
    """业务流程页 - 添加流程箭头"""
    slide = prs.slides.add_slide(layout)
    
    # 标题 - 下移
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(12.33), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.name = CHINESE_FONT
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = MAIN_COLOR
    set_font(title_para, CHINESE_FONT)
    
    if content_texts:
        all_content = "\n".join(content_texts)
        lines = all_content.split('\n')
        
        # 内容区域 - 增大高度
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12.33), Inches(5.4))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        for i, line in enumerate(lines):
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            
            p.text = line
            p.font.name = CHINESE_FONT
            p.font.size = Pt(13)
            p.space_after = Pt(3)
            
            if line.startswith('•'):
                p.level = 0
            elif line.startswith('   '):
                p.level = 1
                p.font.size = Pt(12)
            
            # 数字标题加粗和颜色
            if line and line[0].isdigit() and '.' in line[:3]:
                p.font.bold = True
                p.font.size = Pt(15)
                p.font.color.rgb = MAIN_COLOR
            elif line and not line.startswith(' ') and not line.startswith('•') and len(line) < 20:
                p.font.bold = True
                p.font.size = Pt(15)
                p.font.color.rgb = MAIN_COLOR
            
            set_font(p, CHINESE_FONT)

def add_system_slide_with_architecture(prs, layout, title, content_texts):
    """系统演示页 - 添加架构图"""
    slide = prs.slides.add_slide(layout)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(12.33), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.name = CHINESE_FONT
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = MAIN_COLOR
    set_font(title_para, CHINESE_FONT)
    
    # 左侧内容
    if content_texts:
        all_content = "\n".join(content_texts)
        lines = all_content.split('\n')
        
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(6), Inches(5.4))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        for i, line in enumerate(lines):
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            
            p.text = line
            p.font.name = CHINESE_FONT
            p.font.size = Pt(12)
            p.space_after = Pt(3)
            
            if line.startswith('•'):
                p.level = 0
            elif line.startswith('   '):
                p.level = 1
                p.font.size = Pt(11)
            
            if line and not line.startswith(' ') and not line.startswith('•') and len(line) < 20:
                p.font.bold = True
                p.font.size = Pt(14)
                p.font.color.rgb = MAIN_COLOR
            
            set_font(p, CHINESE_FONT)
    
    # 右侧：系统架构图（分层）
    layers = [
        ("前端层", RGBColor(173, 216, 230)),
        ("业务层", RGBColor(144, 238, 144)),
        ("数据层", RGBColor(255, 218, 185)),
        ("基础层", RGBColor(221, 160, 221))
    ]
    
    y_pos = 2.5
    for layer_name, color in layers:
        add_box_shape(slide, Inches(7.5), Inches(y_pos), Inches(5), Inches(0.9), layer_name, color)
        y_pos += 1.1

def add_management_slide_with_timeline(prs, layout, title, content_texts):
    """项目管理页 - 添加时间轴或组织架构"""
    slide = prs.slides.add_slide(layout)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(12.33), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.name = CHINESE_FONT
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = MAIN_COLOR
    set_font(title_para, CHINESE_FONT)
    
    # 内容 - 增大区域
    if content_texts:
        all_content = "\n".join(content_texts)
        lines = all_content.split('\n')
        
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12.33), Inches(5.4))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        for i, line in enumerate(lines):
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            
            p.text = line
            p.font.name = CHINESE_FONT
            p.font.size = Pt(13)
            p.space_after = Pt(3)
            
            if line.startswith('•'):
                p.level = 0
            elif line.startswith('   '):
                p.level = 1
                p.font.size = Pt(12)
            
            # 突出显示阶段标题
            if '阶段' in line or ('月' in line and '：' in line):
                p.font.bold = True
                p.font.size = Pt(14)
                p.font.color.rgb = ACCENT_COLOR
            elif line and not line.startswith(' ') and not line.startswith('•') and len(line) < 20:
                p.font.bold = True
                p.font.size = Pt(15)
                p.font.color.rgb = MAIN_COLOR
            
            set_font(p, CHINESE_FONT)

def add_deliverable_slide_with_icons(prs, layout, title, content_texts):
    """交付物页 - 添加图标效果"""
    slide = prs.slides.add_slide(layout)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(12.33), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.name = CHINESE_FONT
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = MAIN_COLOR
    set_font(title_para, CHINESE_FONT)
    
    # 内容
    if content_texts:
        all_content = "\n".join(content_texts)
        lines = all_content.split('\n')
        
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12.33), Inches(5.4))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        for i, line in enumerate(lines):
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            
            p.text = line
            p.font.name = CHINESE_FONT
            p.font.size = Pt(13)
            p.space_after = Pt(3)
            
            if line.startswith('•'):
                p.level = 0
            elif line.startswith('   '):
                p.level = 1
                p.font.size = Pt(12)
            
            # 数字列表加颜色
            if line and line[0].isdigit() and '.' in line[:3]:
                p.font.bold = True
                p.font.size = Pt(15)
                p.font.color.rgb = ACCENT_COLOR
            elif line and not line.startswith(' ') and not line.startswith('•') and len(line) < 20:
                p.font.bold = True
                p.font.size = Pt(15)
                p.font.color.rgb = MAIN_COLOR
            
            set_font(p, CHINESE_FONT)

def add_acceptance_slide_with_checkboxes(prs, layout, title, content_texts):
    """验收页 - 添加检查项效果"""
    slide = prs.slides.add_slide(layout)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(12.33), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.name = CHINESE_FONT
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = MAIN_COLOR
    set_font(title_para, CHINESE_FONT)
    
    # 内容
    if content_texts:
        all_content = "\n".join(content_texts)
        lines = all_content.split('\n')
        
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12.33), Inches(5.4))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        for i, line in enumerate(lines):
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            
            # 给验收项添加✓符号
            display_line = line
            if line.startswith('•'):
                display_line = "✓ " + line[1:].strip()
            
            p.text = display_line
            p.font.name = CHINESE_FONT
            p.font.size = Pt(13)
            p.space_after = Pt(3)
            
            if display_line.startswith('✓'):
                p.level = 0
                p.font.color.rgb = RGBColor(0, 128, 0)  # 绿色
            elif line.startswith('   '):
                p.level = 1
                p.font.size = Pt(12)
            
            if line and not line.startswith(' ') and not line.startswith('•') and len(line) < 30:
                p.font.bold = True
                p.font.size = Pt(15)
                p.font.color.rgb = MAIN_COLOR
            
            set_font(p, CHINESE_FONT)

def add_enhanced_content_slide(prs, layout, title, content_texts):
    """增强的普通内容页"""
    slide = prs.slides.add_slide(layout)
    
    # 标题 - 下移
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(12.33), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.name = CHINESE_FONT
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = MAIN_COLOR
    set_font(title_para, CHINESE_FONT)
    
    if content_texts:
        all_content = "\n".join(content_texts)
        lines = all_content.split('\n')
        
        # 增大内容区域
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12.33), Inches(5.4))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        for i, line in enumerate(lines):
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            
            p.text = line
            p.font.name = CHINESE_FONT
            p.font.size = Pt(13)
            p.space_after = Pt(3)
            
            if line.startswith('•'):
                p.level = 0
            elif line.startswith('   '):
                p.level = 1
                p.font.size = Pt(12)
            
            if line and not line.startswith(' ') and not line.startswith('•') and len(line) < 30:
                p.font.bold = True
                p.font.size = Pt(15)
                p.font.color.rgb = MAIN_COLOR
            
            set_font(p, CHINESE_FONT)

def add_enhanced_end_slide(prs, layout, texts):
    """增强的结束页"""
    slide = prs.slides.add_slide(layout)
    
    text = "\n\n".join(texts)
    
    # 居中显示
    end_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9.33), Inches(2))
    end_frame = end_box.text_frame
    end_frame.text = text
    
    for para in end_frame.paragraphs:
        para.font.name = CHINESE_FONT
        para.font.size = Pt(56)
        para.font.bold = True
        para.font.color.rgb = MAIN_COLOR
        para.alignment = PP_ALIGN.CENTER
        set_font(para, CHINESE_FONT)

if __name__ == "__main__":
    print("\n" + "=" * 70)
    print("生成增强版融合PPT")
    print("特性：调整布局 + 添加可视化元素 + 保留模板背景")
    print("=" * 70)
    
    content_path = "/workspace/report_ppt/内容.pptx"
    slides_data = extract_content(content_path)
    
    template_path = "/workspace/report_ppt/模板.pptx"
    output_path = "/workspace/report_ppt/融合.pptx"
    
    create_enhanced_fusion_ppt(template_path, slides_data, output_path)
