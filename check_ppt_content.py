#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation

def check_ppt_content(ppt_path):
    """检查PPT的内容"""
    prs = Presentation(ppt_path)
    
    print("=" * 60)
    print("PPT内容检查：{}".format(ppt_path))
    print("=" * 60)
    print("总幻灯片数：{}".format(len(prs.slides)))
    print("")
    
    slides_list = list(prs.slides)
    for idx, slide in enumerate(slides_list[:10]):  # 只看前10张
        print("--- 幻灯片 {} ---".format(idx + 1))
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_preview = shape.text[:100].replace('\n', ' ')
                print("  {}".format(text_preview))
        print("")

if __name__ == "__main__":
    check_ppt_content("/workspace/report_ppt/Wednesday_Report.pptx")
