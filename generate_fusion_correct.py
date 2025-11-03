#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement

CHINESE_FONT = "WenQuanYi Zen Hei"
MAIN_COLOR = RGBColor(0, 70, 160)

def set_font(paragraph, font_name):
    """设置段落字体，包括东亚字体"""
    try:
        for run in paragraph.runs:
            run.font.name = font_name
            if hasattr(run.font, '_element'):
                rPr = run.font._element.get_or_add_rPr()
                ea = OxmlElement('a:ea')
                ea.set('typeface', font_name)
                rPr.append(ea)
                cs = OxmlElement('a:cs')
                cs.set('typeface', font_name)
                rPr.append(cs)
    except:
        pass

def extract_content(content_ppt_path):
    """提取内容PPT的所有数据"""
    print("\n步骤1: 读取内容.pptx")
    print("-" * 70)
    prs = Presentation(content_ppt_path)
    
    slides_data = []
    slides_list = list(prs.slides)
    
    for idx, slide in enumerate(slides_list):
        slide_info = {
            'index': idx + 1,
            'texts': []
        }
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_info['texts'].append(shape.text.strip())
        
        slides_data.append(slide_info)
        
        # 显示提取的内容
        title = slide_info['texts'][0] if slide_info['texts'] else "(空)"
        print("  幻灯片 {}: {}".format(idx + 1, title[:40]))
    
    print("\n✓ 共提取 {} 张幻灯片的内容\n".format(len(slides_data)))
    return slides_data

def create_fusion_ppt(template_path, slides_data, output_path):
    """创建融合PPT - 只保留模板样式，内容完全来自内容.pptx"""
    print("步骤2: 加载模板.pptx（只用于获取布局）")
    print("-" * 70)
    
    # 加载模板，只为了获取布局样式
    template_prs = Presentation(template_path)
    
    # 复制模板的幻灯片尺寸和布局
    print("  提取模板布局...")
    title_layout = template_prs.slide_layouts[12]  # 1_标题幻灯片
    section_layout = template_prs.slide_layouts[13]  # 1_节标题
    blank_layout = template_prs.slide_layouts[6]  # 空白
    
    # 创建全新的演示文稿，复制模板的母版
    print("  创建新演示文稿（保留模板样式）...")
    prs = Presentation(template_path)
    
    # 删除模板中的所有幻灯片（只保留母版和布局）
    print("  删除模板中的所有幻灯片...")
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]
    
    print("✓ 模板已清空，只保留样式\n")
    
    print("步骤3: 使用内容.pptx的内容填充")
    print("-" * 70)
    
    # 重新获取布局（因为prs已经是新的）
    title_layout = prs.slide_layouts[12]
    section_layout = prs.slide_layouts[13]
    blank_layout = prs.slide_layouts[6]
    
    for slide_data in slides_data:
        idx = slide_data['index']
        texts = slide_data['texts']
        
        if not texts:
            continue
        
        first_text = texts[0]
        
        # 标题页
        if idx == 1:
            add_title_slide(prs, title_layout, texts)
            print("  [{}] 标题页: {}".format(idx, first_text[:30]))
        
        # 节标题或内容页
        elif first_text.startswith(("一、", "二、", "三、", "四、", "五、", "六、")):
            section_title = first_text[2:]
            add_content_slide(prs, blank_layout, first_text, texts[1:])
            print("  [{}] 内容页: {}".format(idx, first_text[:30]))
        
        # 结束页
        elif "谢谢" in first_text or "Q & A" in first_text:
            add_end_slide(prs, blank_layout, texts)
            print("  [{}] 结束页".format(idx))
        
        # 普通内容页
        else:
            add_content_slide(prs, blank_layout, first_text, texts[1:] if len(texts) > 1 else [])
            print("  [{}] 内容页: {}".format(idx, first_text[:30]))
    
    print("\n✓ 内容填充完成\n")
    
    print("步骤4: 保存融合.pptx")
    print("-" * 70)
    prs.save(output_path)
    print("✓ 文件已保存: {}\n".format(output_path))
    
    print("=" * 70)
    print("✅ 融合PPT生成成功！")
    print("=" * 70)
    print("📄 文件: {}".format(output_path))
    print("📊 幻灯片: {} 张".format(len(slides_data)))
    print("🎨 样式来源: 模板.pptx")
    print("📝 内容来源: 内容.pptx")
    print("=" * 70)

def add_title_slide(prs, layout, texts):
    """添加标题页"""
    slide = prs.slides.add_slide(layout)
    
    title = texts[0] if len(texts) > 0 else "项目汇报"
    subtitle = texts[1] if len(texts) > 1 else ""
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.name = CHINESE_FONT
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = MAIN_COLOR
    title_para.alignment = PP_ALIGN.CENTER
    set_font(title_para, CHINESE_FONT)
    
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11.33), Inches(0.8))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.name = CHINESE_FONT
        subtitle_para.font.size = Pt(28)
        subtitle_para.font.color.rgb = RGBColor(89, 89, 89)
        subtitle_para.alignment = PP_ALIGN.CENTER
        set_font(subtitle_para, CHINESE_FONT)

def add_content_slide(prs, layout, title, content_texts):
    """添加内容页"""
    slide = prs.slides.add_slide(layout)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.name = CHINESE_FONT
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = MAIN_COLOR
    set_font(title_para, CHINESE_FONT)
    
    if content_texts:
        all_content = "\n".join(content_texts)
        lines = all_content.split('\n')
        
        # 智能调整字体大小
        total_lines = len(lines)
        if total_lines > 50:
            base_font_size = 11
            title_font_size = 13
            space_after = 2
        elif total_lines > 35:
            base_font_size = 12
            title_font_size = 14
            space_after = 3
        elif total_lines > 25:
            base_font_size = 13
            title_font_size = 15
            space_after = 3
        else:
            base_font_size = 14
            title_font_size = 16
            space_after = 4
        
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.33), Inches(5.8))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        for i, line in enumerate(lines):
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            
            p.text = line
            p.font.name = CHINESE_FONT
            p.font.size = Pt(base_font_size)
            p.space_after = Pt(space_after)
            
            if line.startswith('•'):
                p.level = 0
            elif line.startswith('   •'):
                p.level = 1
                p.font.size = Pt(base_font_size - 1)
            elif line.startswith('   '):
                p.level = 1
                p.font.size = Pt(base_font_size - 1)
            
            # 加粗小标题
            if line and not line.startswith(' ') and not line.startswith('•') and \
               len(line) < 30 and ':' not in line and '：' not in line and \
               not line.startswith(('第', '已', '当', '下', '验收', '潜在')):
                p.font.bold = True
                p.font.size = Pt(title_font_size)
                p.font.color.rgb = MAIN_COLOR
            
            set_font(p, CHINESE_FONT)

def add_end_slide(prs, layout, texts):
    """添加结束页"""
    slide = prs.slides.add_slide(layout)
    
    text = "\n\n".join(texts)
    
    end_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9.33), Inches(2))
    end_frame = end_box.text_frame
    end_frame.text = text
    
    for para in end_frame.paragraphs:
        para.font.name = CHINESE_FONT
        para.font.size = Pt(52)
        para.font.bold = True
        para.font.color.rgb = MAIN_COLOR
        para.alignment = PP_ALIGN.CENTER
        set_font(para, CHINESE_FONT)

if __name__ == "__main__":
    print("\n" + "=" * 70)
    print("开始生成融合PPT")
    print("说明：使用内容.pptx的内容 + 模板.pptx的样式")
    print("=" * 70)
    
    # 提取内容
    content_path = "/workspace/report_ppt/内容.pptx"
    slides_data = extract_content(content_path)
    
    # 创建融合PPT
    template_path = "/workspace/report_ppt/模板.pptx"
    output_path = "/workspace/report_ppt/融合.pptx"
    
    create_fusion_ppt(template_path, slides_data, output_path)
