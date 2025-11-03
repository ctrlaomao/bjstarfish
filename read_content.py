#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt

def read_content_ppt(ppt_path):
    """读取内容PPT的详细信息"""
    print("=" * 70)
    print("内容PPT分析：{}".format(ppt_path))
    print("=" * 70)
    
    prs = Presentation(ppt_path)
    
    # 1. 基本信息
    print("\n【1. 基本信息】")
    print("幻灯片尺寸：宽 {} x 高 {}".format(
        prs.slide_width / Inches(1), 
        prs.slide_height / Inches(1)
    ))
    print("总幻灯片数：{}".format(len(prs.slides)))
    
    # 2. 提取所有幻灯片内容
    print("\n【2. 详细内容】")
    slides_list = list(prs.slides)
    
    for idx, slide in enumerate(slides_list):
        print("\n" + "=" * 70)
        print("幻灯片 {} / {}".format(idx + 1, len(slides_list)))
        print("=" * 70)
        print("布局：{}".format(slide.slide_layout.name))
        
        # 提取所有文本
        all_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                all_texts.append(shape.text.strip())
        
        if all_texts:
            for text_idx, text in enumerate(all_texts):
                print("\n[文本块 {}]".format(text_idx + 1))
                print(text)
        else:
            print("\n(无文本内容)")
    
    print("\n" + "=" * 70)
    print("内容分析完成！")
    print("=" * 70)
    
    return prs

if __name__ == "__main__":
    content_path = "/workspace/report_ppt/内容.pptx"
    read_content_ppt(content_path)
