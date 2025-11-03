#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt

def read_template_ppt(ppt_path):
    """读取模板PPT的详细信息"""
    print("=" * 70)
    print("模板PPT内容分析：{}".format(ppt_path))
    print("=" * 70)
    
    prs = Presentation(ppt_path)
    
    # 1. 基本信息
    print("\n【1. 基本信息】")
    print("幻灯片尺寸：宽 {} x 高 {}".format(
        prs.slide_width / Inches(1), 
        prs.slide_height / Inches(1)
    ))
    print("总幻灯片数：{}".format(len(prs.slides)))
    print("可用布局数：{}".format(len(prs.slide_layouts)))
    
    # 2. 布局信息
    print("\n【2. 幻灯片布局】")
    for idx, layout in enumerate(prs.slide_layouts):
        print("  布局 {}: {} (占位符: {})".format(
            idx, layout.name, len(layout.placeholders)
        ))
    
    # 3. 前5张幻灯片内容预览
    print("\n【3. 幻灯片内容预览（前5张）】")
    slides_list = list(prs.slides)
    
    for idx, slide in enumerate(slides_list[:5]):
        print("\n--- 幻灯片 {} ---".format(idx + 1))
        print("使用的布局：{}".format(slide.slide_layout.name))
        
        # 提取文本内容
        text_count = 0
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                text_preview = text[:80] + "..." if len(text) > 80 else text
                print("  [文本 {}] {}".format(text_count + 1, text_preview))
                text_count += 1
                
                # 显示字体信息（如果有）
                if hasattr(shape, "text_frame") and len(shape.text_frame.paragraphs) > 0:
                    para = shape.text_frame.paragraphs[0]
                    if len(para.runs) > 0:
                        run = para.runs[0]
                        font_info = "    字体: {}, 大小: {}, 加粗: {}".format(
                            run.font.name if run.font.name else "默认",
                            run.font.size if run.font.size else "默认",
                            run.font.bold if run.font.bold is not None else "默认"
                        )
                        print(font_info)
        
        if text_count == 0:
            print("  (无文本内容)")
    
    # 4. 统计所有幻灯片
    print("\n【4. 完整幻灯片列表】")
    for idx, slide in enumerate(slides_list):
        titles = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                if len(text) < 50:  # 简短文本可能是标题
                    titles.append(text)
        
        title_str = " | ".join(titles[:2]) if titles else "(无标题)"
        print("  幻灯片 {}: {}".format(idx + 1, title_str[:60]))
    
    print("\n" + "=" * 70)
    print("模板分析完成！")
    print("=" * 70)
    
    return prs

if __name__ == "__main__":
    template_path = "/workspace/report_ppt/模板.pptx"
    read_template_ppt(template_path)
